"""
Azure Consumption Plan – BCG (Boston Consulting Group) Style PowerPoint
Universitas Terbuka | April 2026 – March 2027
─────────────────────────────────────────────────────────────────────────
BCG Style: Green brand color, bold takeaway bars, geometric accents,
"So what?" callouts, clean data-heavy layouts, Calibri/Helvetica fonts,
traffic-light indicators, white background, minimal decoration.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── BCG Color Palette ──────────────────────────────────────────────────
BCG_GREEN = RGBColor(0x00, 0xA6, 0x51)
BCG_DARK_GREEN = RGBColor(0x00, 0x6B, 0x3F)
BCG_LIGHT_GREEN = RGBColor(0x7E, 0xD3, 0x57)
BCG_FOREST = RGBColor(0x1B, 0x4D, 0x3E)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY = RGBColor(0x58, 0x58, 0x58)
LIGHT_GRAY = RGBColor(0x9B, 0x9B, 0x9B)
VERY_LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BCG_BLUE = RGBColor(0x00, 0x6F, 0xAF)
BCG_TEAL = RGBColor(0x00, 0x97, 0x93)
BCG_ORANGE = RGBColor(0xED, 0x8B, 0x00)
BCG_RED = RGBColor(0xC8, 0x10, 0x2E)
BCG_PURPLE = RGBColor(0x6C, 0x20, 0x7E)
AMBER = RGBColor(0xF5, 0xA6, 0x23)

CHART_PALETTE = [
    BCG_GREEN, BCG_BLUE, BCG_TEAL, BCG_ORANGE,
    BCG_PURPLE, BCG_LIGHT_GREEN, LIGHT_GRAY,
    RGBColor(0x00, 0x4E, 0x7C), RGBColor(0xA0, 0xD9, 0x68),
    RGBColor(0xE0, 0x60, 0x60),
]

FONT = "Calibri"

# ── Data ───────────────────────────────────────────────────────────────
MONTHS_FY26 = ["Jul-25","Aug-25","Sep-25","Oct-25","Nov-25","Dec-25",
               "Jan-26","Feb-26","Mar-26","Apr-26"]
SUBS = {
    "AZURE APP UT":             [0, 0, 0, 295, 377, 425, 420, 401, 574, 205],
    "AZURE E-LEARNING":         [1231, 1232, 4071, 45894, 44684, 17745, 14877, 46491, 90451, 29456],
    "AZURE IN OPEN":            [99988, 122782, 92467, 120836, 140956, 131736, 135660, 109095, 84438, 0],
    "AZURE SRS":                [0, 0, 0, 0, 0, 0, 4481, 8265, 9280, 3238],
    "Azure subscription 1":     [4051, 4112, 3943, 4166, 3601, 3675, 3761, 3467, 3763, 1058],
    "AZURE THE":                [0, 0, 80, 125, 122, 125, 125, 114, 125, 44],
    "DSI ELEARNING RND":        [0, 0, 157, 379, 557, 815, 1161, 756, 539, 155],
    "Microsoft Azure":          [2202, 2178, 2091, 2154, 2154, 2158, 2107, 1270, 1191, 419],
    "SUB-PLATFORM-NETWORK-IDC": [0, 0, 0, 0, 0, 0, 0, 251, 414, 145],
    "UNKNOWN":                  [105, 190, 243, 254, 262, 286, 182, 266, 192, 0],
}
GROWTH = 0.20

RI_SERVICES = [
    ("AKS Cluster", "D4s_v5 ×3 nodes", 620, 400, 3),
    ("Frontend (Vue.js)", "P1v3 App Service", 138, 90, 2),
    ("Backend (PHP)", "D2s_v5 Container", 310, 200, 2),
    ("PostgreSQL Master", "GP_D4s_v3", 530, 345, 1),
    ("PostgreSQL Slave", "GP_D2s_v3", 265, 172, 1),
    ("DB Storage", "256 GB SSD", 80, 80, 1),
    ("Azure CDN", "10 TB/mo", 850, 850, 1),
    ("Load Balancer", "Standard LB", 25, 25, 1),
    ("Monitoring", "Log Analytics", 60, 60, 1),
]

AI_SERVICES = [
    ("Azure OpenAI (GPT-4o)", "1M tokens/day", 4500),
    ("Azure OpenAI (Embeddings)", "ada-002", 800),
    ("Azure AI Search", "S1 Standard", 750),
    ("Cognitive Svc (Speech)", "TTS", 400),
    ("Blob Storage", "Hot, 5 TB", 120),
    ("Cosmos DB", "4000 RU/s", 580),
    ("Container Apps", "2 instances", 450),
    ("Redis Cache", "C2 Standard", 170),
    ("App Service (Portal)", "P2v3", 280),
    ("DevOps / CI-CD", "Agents", 150),
]

# ── Computed ──────────────────────────────────────────────────────────
fy26_actual = sum(sum(v) for v in SUBS.values())
fy26_9mo = sum(sum(v[:9]) for v in SUBS.values())
fy26_annualized = round(fy26_9mo / 9 * 12)
fy26_monthly = [sum(SUBS[s][m] for s in SUBS) for m in range(10)]
fy26_avg = sum(fy26_monthly) / len(fy26_monthly)

FY26_IDX = [9, None, None, 0, 1, 2, 3, 4, 5, 6, 7, 8]
projected = []
for mi in range(12):
    idx = FY26_IDX[mi]
    base = fy26_avg if idx is None else fy26_monthly[idx]
    projected.append(round(base * (1 + GROWTH)))
plan_annual = sum(projected)

ri_payg_r = sum(s[2]*s[4] for s in RI_SERVICES)
ri_ri_r = sum(s[3]*s[4] for s in RI_SERVICES)
ri_annual_4 = ri_ri_r * 4 * 12
ri_savings = (ri_payg_r - ri_ri_r) * 4 * 12
ai_mo = sum(s[2] for s in AI_SERVICES)
ai_annual = ai_mo * 12
grand_total = plan_annual + ri_annual_4 + ai_annual

sub_proj = {}
for s, v in SUBS.items():
    avg = sum(v)/len(v); t = 0
    for mi in range(12):
        idx = FY26_IDX[mi]
        t += (avg if idx is None else v[idx]) * (1 + GROWTH)
    sub_proj[s] = round(t)

def fmt(v):
    if v >= 1_000_000: return f"${v/1_000_000:.1f}M"
    if v >= 1_000: return f"${v/1_000:.0f}K"
    return f"${v:,.0f}"

def fmtf(v): return f"${v:,.0f}"

# ── Presentation ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide_num = [0]

def add_rect(sl, l, t, w, h, c, border=None):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c
    if border: s.line.color.rgb = border; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def add_txt(sl, l, t, w, h, txt, sz=12, bold=False, c=DARK_GRAY,
            al=PP_ALIGN.LEFT, fn=FONT, italic=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.bold = bold; p.font.italic = italic
    p.font.color.rgb = c; p.font.name = fn; p.alignment = al
    return tb

def add_multiline(sl, l, t, w, h, lines, sz=10, c=DARK_GRAY, bullet="▪"):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {line}" if bullet else line
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = FONT
        p.space_after = Pt(2)
    return tb

def bcg_takeaway(sl, text):
    """BCG-style green takeaway bar at top of slide."""
    slide_num[0] += 1
    add_rect(sl, 0, 0, Inches(13.333), Inches(0.95), BCG_GREEN)
    add_txt(sl, Inches(0.6), Inches(0.15), Inches(11.5), Inches(0.7),
            text, sz=18, bold=True, c=WHITE)
    # Slide number
    add_txt(sl, Inches(12.3), Inches(0.2), Inches(0.8), Inches(0.4),
            str(slide_num[0]), sz=14, bold=True, c=WHITE, al=PP_ALIGN.RIGHT)

def bcg_footer(sl, src="Source: Azure ACR Data FY26; Azure Pricing Calculator; Internal estimates"):
    add_rect(sl, 0, Inches(7.15), Inches(13.333), Inches(0.35), BG_LIGHT)
    add_txt(sl, Inches(0.6), Inches(7.18), Inches(10), Inches(0.25),
            src, sz=7, c=LIGHT_GRAY, italic=True)
    add_txt(sl, Inches(11.0), Inches(7.18), Inches(2), Inches(0.25),
            "CONFIDENTIAL", sz=7, bold=True, c=LIGHT_GRAY, al=PP_ALIGN.RIGHT)

def so_what_box(sl, x, y, w, h, text):
    """BCG 'So what?' callout box."""
    add_rect(sl, x, y, w, Inches(0.3), BCG_GREEN)
    add_txt(sl, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.25),
            "SO WHAT?", sz=9, bold=True, c=WHITE)
    box = add_rect(sl, x, y + Inches(0.3), w, h - Inches(0.3), BG_LIGHT, border=VERY_LIGHT_GRAY)
    add_txt(sl, x + Inches(0.15), y + Inches(0.35), w - Inches(0.3), h - Inches(0.4),
            text, sz=10, c=DARK_GRAY)
    return box

def bcg_kpi(sl, x, y, w, h, value, label, color=BCG_GREEN):
    """BCG-style metric card with thick left bar."""
    add_rect(sl, x, y, w, h, WHITE, border=VERY_LIGHT_GRAY)
    add_rect(sl, x, y, Pt(6), h, color)  # left accent bar
    add_txt(sl, x + Inches(0.15), y + Inches(0.08), w - Inches(0.2), Inches(0.5),
            value, sz=24, bold=True, c=color, al=PP_ALIGN.LEFT)
    add_txt(sl, x + Inches(0.15), y + Inches(0.55), w - Inches(0.2), Inches(0.25),
            label, sz=9, c=MID_GRAY)

def bcg_table(sl, data, l, t, cw, hdr_bg=BCG_GREEN):
    rows = len(data); cols = len(data[0])
    ts = sl.shapes.add_table(rows, cols, l, t, sum(cw), Inches(0.3 * rows))
    tbl = ts.table
    for ci, w in enumerate(cw): tbl.columns[ci].width = w
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = str(val)
            is_hdr = ri == 0; is_total = ri == rows - 1
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT; p.font.size = Pt(9 if not is_hdr else 10)
                p.font.bold = is_hdr or is_total
                p.font.color.rgb = WHITE if is_hdr else (BCG_DARK_GREEN if is_total else DARK_GRAY)
                p.alignment = PP_ALIGN.RIGHT if ci > 0 else PP_ALIGN.LEFT
            if is_hdr:
                cell.fill.solid(); cell.fill.fore_color.rgb = hdr_bg
            elif is_total:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
            elif ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = BG_LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    return tbl

def traffic_light(sl, x, y, color, text):
    """Traffic light indicator dot + text."""
    add_rect(sl, x, y + Inches(0.04), Inches(0.15), Inches(0.15), color)
    add_txt(sl, x + Inches(0.2), y, Inches(2.5), Inches(0.22), text, sz=9, c=DARK_GRAY)

# =====================================================================
# SLIDE 1 – COVER (BCG style: white bg, large green accent, bold type)
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
# Left green accent panel
add_rect(sl, 0, 0, Inches(0.4), Inches(7.5), BCG_GREEN)
# Top green stripe
add_rect(sl, 0, 0, Inches(13.333), Inches(0.08), BCG_GREEN)
# Bottom green stripe
add_rect(sl, 0, Inches(7.42), Inches(13.333), Inches(0.08), BCG_GREEN)

add_txt(sl, Inches(1.0), Inches(0.6), Inches(4), Inches(0.4),
        "UNIVERSITAS TERBUKA", sz=13, bold=True, c=BCG_GREEN)

add_rect(sl, Inches(1.0), Inches(1.1), Inches(3), Pt(3), BCG_GREEN)

add_txt(sl, Inches(1.0), Inches(1.5), Inches(10), Inches(1.2),
        "Azure Cloud\nConsumption Plan", sz=44, bold=True, c=BLACK)

add_txt(sl, Inches(1.0), Inches(3.5), Inches(10), Inches(0.6),
        "April 2026 – March 2027  |  12-Month Rolling Plan", sz=20, c=MID_GRAY)

add_txt(sl, Inches(1.0), Inches(4.3), Inches(10), Inches(0.4),
        "Growth Rate: 20%  |  Prepared: April 2026  |  Currency: USD", sz=14, c=LIGHT_GRAY)

# KPI strip
kpi_strip_y = Inches(5.2)
add_rect(sl, Inches(1.0), kpi_strip_y, Inches(11.5), Inches(1.6), BG_LIGHT, border=VERY_LIGHT_GRAY)
kpis = [
    (fmtf(fy26_actual), "FY26 Actual\n(10 months)", LIGHT_GRAY),
    (fmtf(fy26_annualized), "FY26 Annualized\n(12 months)", MID_GRAY),
    (fmtf(plan_annual), "Existing Subs\nProjected", BCG_BLUE),
    (fmtf(ai_annual), "AI Course Gen\n(New)", BCG_TEAL),
    (fmtf(ri_annual_4), "Ujian Online RI\n(4 Regions)", BCG_ORANGE),
    (fmtf(grand_total), "Grand Total\n(12 months)", BCG_GREEN),
]
x = Inches(1.2)
for val, label, clr in kpis:
    add_rect(sl, x, kpi_strip_y + Inches(0.1), Inches(0.06), Inches(1.4), clr)  # left bar
    add_txt(sl, x + Inches(0.15), kpi_strip_y + Inches(0.15), Inches(1.6), Inches(0.5),
            val, sz=18, bold=True, c=clr)
    add_txt(sl, x + Inches(0.15), kpi_strip_y + Inches(0.7), Inches(1.6), Inches(0.6),
            label, sz=9, c=MID_GRAY)
    x += Inches(1.9)

add_txt(sl, Inches(10.5), Inches(7.05), Inches(2.5), Inches(0.3),
        "CONFIDENTIAL", sz=9, bold=True, c=LIGHT_GRAY, al=PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 2 – EXECUTIVE SUMMARY
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bcg_takeaway(sl,
    f"Universitas Terbuka cloud investment plan totals {fmtf(grand_total)} over 12 months, "
    "driven by 20% organic growth and two strategic initiatives")

# Three pillars
pillars = [
    ("CURRENT STATE", BCG_BLUE, [
        f"FY26 spend: {fmtf(fy26_actual)} (10 months actual)",
        f"Annualized baseline: {fmtf(fy26_annualized)}",
        "10 active subscriptions across e-learning & platform",
        "AZURE IN OPEN dominant (73%) but declining",
        "AZURE E-LEARNING surging ($1K → $90K/mo)",
    ]),
    ("GROWTH INITIATIVES", BCG_GREEN, [
        f"AI Course Generation: {fmtf(ai_annual)}/year",
        "GPT-4o, AI Search, Speech, Cosmos DB",
        f"Ujian Online Platform (4 regions): {fmtf(ri_annual_4)}/year",
        "AKS, Vue.js, PHP, PostgreSQL M-S per region",
        f"RI savings vs PAYG: {fmtf(ri_savings)}/year",
    ]),
    ("TOTAL INVESTMENT", BCG_ORANGE, [
        f"Existing subscriptions (20%): {fmtf(plan_annual)}",
        f"AI Course Generation: {fmtf(ai_annual)}",
        f"Ujian Online RI (4 regions): {fmtf(ri_annual_4)}",
        f"Grand Total: {fmtf(grand_total)}",
        "Plan period: April 2026 – March 2027",
    ]),
]

for i, (title, color, bullets) in enumerate(pillars):
    x = Inches(0.6) + i * Inches(4.2)
    w = Inches(3.9)
    add_rect(sl, x, Inches(1.2), w, Inches(0.4), color)
    add_txt(sl, x + Inches(0.1), Inches(1.23), w, Inches(0.35),
            title, sz=12, bold=True, c=WHITE)
    add_multiline(sl, x + Inches(0.1), Inches(1.75), w - Inches(0.2), Inches(3.0),
                  bullets, sz=10, c=DARK_GRAY)

# Bottom KPI cards
x = Inches(0.6)
for val, label, clr in [
    (fmtf(fy26_annualized), "FY26 Baseline", LIGHT_GRAY),
    (fmtf(grand_total), "Plan Total", BCG_GREEN),
    ("20%", "Growth Rate", BCG_BLUE),
    (fmtf(ri_savings), "RI Savings/yr", BCG_TEAL),
]:
    bcg_kpi(sl, x, Inches(5.5), Inches(2.85), Inches(0.85), val, label, clr)
    x += Inches(3.1)

bcg_footer(sl)

# =====================================================================
# SLIDE 3 – FY26 CURRENT STATE ANALYSIS
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bcg_takeaway(sl,
    f"FY26 Azure consumption reached {fmtf(fy26_actual)} across 10 months, "
    "with AZURE IN OPEN accounting for 73% of total spend")

# Bar chart
cd = CategoryChartData()
cd.categories = MONTHS_FY26
cd.add_series("Monthly Consumption (USD)", fy26_monthly)
cf = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.6), Inches(1.3), Inches(7.8), Inches(4.8), cd)
ch = cf.chart; ch.has_legend = False; ch.style = 2
pl = ch.plots[0]; pl.gap_width = 50
sr = pl.series[0]; sr.format.fill.solid(); sr.format.fill.fore_color.rgb = BCG_GREEN
sr.has_data_labels = True
dl = sr.data_labels; dl.font.size = Pt(8); dl.font.color.rgb = MID_GRAY
dl.number_format = '$#,##0'; dl.show_value = True
ch.value_axis.major_gridlines.format.line.color.rgb = VERY_LIGHT_GRAY
ch.value_axis.format.line.fill.background()
ch.category_axis.format.line.color.rgb = VERY_LIGHT_GRAY
ch.value_axis.tick_labels.font.size = Pt(8); ch.value_axis.tick_labels.number_format = '$#,##0'
ch.category_axis.tick_labels.font.size = Pt(8)

# Right panel: subscription table
top_subs = sorted(SUBS.items(), key=lambda x: sum(x[1]), reverse=True)
td = [["Subscription", "10-mo Total", "Share"]]
for name, vals in top_subs:
    total = sum(vals)
    td.append([name, fmtf(total), f"{total/fy26_actual:.1%}"])
td.append(["TOTAL", fmtf(fy26_actual), "100%"])
bcg_table(sl, td, Inches(8.7), Inches(1.3), [Inches(2.2), Inches(1.2), Inches(0.8)])

# So what box
so_what_box(sl, Inches(8.7), Inches(5.2), Inches(4.2), Inches(1.5),
    "AZURE IN OPEN dominates at 73% but dropped to $0 in Apr-26, "
    "while AZURE E-LEARNING surged 73× from $1.2K to $90K/mo. "
    "This signals a major platform transition requiring proactive capacity planning.")

bcg_footer(sl)

# =====================================================================
# SLIDE 4 – PROJECTED CONSUMPTION
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bcg_takeaway(sl,
    f"Applying 20% growth per-month yields {fmtf(plan_annual)} in projected existing subscription spend "
    f"— an incremental {fmtf(plan_annual - fy26_annualized)} vs FY26 annualized baseline")

# Side-by-side bars
compare_mo = ["Jul","Aug","Sep","Oct","Nov","Dec","Jan","Feb","Mar","Apr"]
proj_compare = [projected[3],projected[4],projected[5],projected[6],projected[7],
                projected[8],projected[9],projected[10],projected[11],projected[0]]
cd = CategoryChartData()
cd.categories = compare_mo
cd.add_series("FY26 Actual", fy26_monthly)
cd.add_series("Projected (20% Growth)", proj_compare)
cf = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.6), Inches(1.3), Inches(8.5), Inches(4.2), cd)
ch = cf.chart; ch.style = 2
pl = ch.plots[0]; pl.gap_width = 40; pl.overlap = -15
pl.series[0].format.fill.solid(); pl.series[0].format.fill.fore_color.rgb = VERY_LIGHT_GRAY
pl.series[1].format.fill.solid(); pl.series[1].format.fill.fore_color.rgb = BCG_GREEN
ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.TOP
ch.legend.include_in_layout = False; ch.legend.font.size = Pt(9)
ch.value_axis.major_gridlines.format.line.color.rgb = VERY_LIGHT_GRAY
ch.value_axis.format.line.fill.background()
ch.value_axis.tick_labels.font.size = Pt(8); ch.value_axis.tick_labels.number_format = '$#,##0'
ch.category_axis.tick_labels.font.size = Pt(9)

# Right KPIs
for i, (val, label, clr) in enumerate([
    (fmtf(fy26_annualized), "FY26 Annualized", LIGHT_GRAY),
    (fmtf(plan_annual), "Plan Annual", BCG_GREEN),
    ("20.0%", "Growth Rate", BCG_BLUE),
    (fmtf(plan_annual - fy26_annualized), "Incremental Spend", BCG_ORANGE),
]):
    bcg_kpi(sl, Inches(9.5), Inches(1.3) + i * Inches(1.0), Inches(3.5), Inches(0.85), val, label, clr)

# Formula box
add_rect(sl, Inches(0.6), Inches(5.7), Inches(8.5), Inches(0.35), BCG_DARK_GREEN)
add_txt(sl, Inches(0.7), Inches(5.72), Inches(8.3), Inches(0.3),
        "FORMULA", sz=10, bold=True, c=WHITE)
add_rect(sl, Inches(0.6), Inches(6.05), Inches(8.5), Inches(0.6), BG_LIGHT, border=VERY_LIGHT_GRAY)
add_txt(sl, Inches(0.8), Inches(6.1), Inches(8.1), Inches(0.5),
        "Apr-26 = FY26 Apr × 1.20   |   May/Jun-26 = AVG(FY26 all months) × 1.20   |   "
        "Jul-26 → Mar-27 = Same calendar month in FY26 × 1.20", sz=9, c=MID_GRAY)

# Traffic lights
traffic_light(sl, Inches(9.5), Inches(5.7), BCG_GREEN, "Growth within budget range")
traffic_light(sl, Inches(9.5), Inches(5.95), AMBER, "AZURE IN OPEN decline risk")
traffic_light(sl, Inches(9.5), Inches(6.2), BCG_GREEN, "E-LEARNING ramp on track")

bcg_footer(sl)

# =====================================================================
# SLIDE 5 – SUBSCRIPTION MIX
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bcg_takeaway(sl,
    "AZURE IN OPEN and AZURE E-LEARNING represent 93% of projected existing subscription spend "
    "— portfolio concentration creates both efficiency and risk")

sorted_subs = sorted(sub_proj.items(), key=lambda x: x[1], reverse=True)
cd = CategoryChartData()
cd.categories = [s[0] for s in sorted_subs]
cd.add_series("Projected Annual (USD)", [s[1] for s in sorted_subs])
cf = sl.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.6), Inches(1.3), Inches(7.5), Inches(5.2), cd)
ch = cf.chart; ch.has_legend = False; ch.style = 2
pl = ch.plots[0]; pl.gap_width = 60
sr = pl.series[0]; sr.format.fill.solid(); sr.format.fill.fore_color.rgb = BCG_GREEN
sr.has_data_labels = True
dl = sr.data_labels; dl.font.size = Pt(9); dl.font.color.rgb = MID_GRAY
dl.number_format = '$#,##0'; dl.show_value = True
ch.value_axis.major_gridlines.format.line.color.rgb = VERY_LIGHT_GRAY
ch.value_axis.format.line.fill.background()
ch.value_axis.tick_labels.font.size = Pt(8); ch.value_axis.tick_labels.number_format = '$#,##0'
ch.category_axis.tick_labels.font.size = Pt(9)

td = [["Subscription", "Annual (USD)", "Share"]]
for name, val in sorted_subs:
    td.append([name, fmtf(val), f"{val/plan_annual:.1%}"])
td.append(["TOTAL", fmtf(plan_annual), "100%"])
bcg_table(sl, td, Inches(8.5), Inches(1.3), [Inches(2.4), Inches(1.2), Inches(0.7)])

so_what_box(sl, Inches(8.5), Inches(5.3), Inches(4.2), Inches(1.2),
    "High concentration in 2 subscriptions (93%). "
    "Consider workload-level cost tagging and consumption alerts to manage risk.")

bcg_footer(sl)

# =====================================================================
# SLIDE 6 – AI COURSE GENERATION
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bcg_takeaway(sl,
    f"AI Course Generation requires {fmtf(ai_annual)}/year in new Azure services — "
    "enabling automated content creation at scale for 300K+ students")

add_rect(sl, Inches(0.6), Inches(1.2), Inches(6), Inches(0.35), BCG_TEAL)
add_txt(sl, Inches(0.7), Inches(1.22), Inches(5.8), Inches(0.3),
        "SERVICE BREAKDOWN", sz=10, bold=True, c=WHITE)

td = [["Azure Service", "SKU", "Monthly (USD)", "Annual (USD)"]]
for svc, desc, mo in AI_SERVICES:
    td.append([svc, desc, fmtf(mo), fmtf(mo * 12)])
td.append(["TOTAL", "", fmtf(ai_mo), fmtf(ai_annual)])
bcg_table(sl, td, Inches(0.6), Inches(1.6),
    [Inches(2.5), Inches(1.5), Inches(1.0), Inches(1.0)], hdr_bg=BCG_TEAL)

# Pie
cd = CategoryChartData()
cd.categories = [s[0] for s in AI_SERVICES]
cd.add_series("Cost", [s[2] for s in AI_SERVICES])
cf = sl.shapes.add_chart(XL_CHART_TYPE.PIE,
    Inches(7.0), Inches(1.3), Inches(3.5), Inches(3.0), cd)
ch = cf.chart; ch.has_legend = False
for i2, pt in enumerate(ch.plots[0].series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = CHART_PALETTE[i2 % len(CHART_PALETTE)]
ch.plots[0].series[0].has_data_labels = True
dl2 = ch.plots[0].series[0].data_labels
dl2.font.size = Pt(7); dl2.show_percentage = True; dl2.show_category_name = True
dl2.show_value = False; dl2.separator = "\n"

# Value prop
add_rect(sl, Inches(7.0), Inches(4.5), Inches(6), Inches(0.35), BCG_TEAL)
add_txt(sl, Inches(7.1), Inches(4.52), Inches(5.8), Inches(0.3),
        "VALUE PROPOSITION", sz=10, bold=True, c=WHITE)
vp = [
    "Automates course creation — reduces manual effort by ~60%",
    "GPT-4o generates lectures, quizzes, and assessments",
    "AI Search enables intelligent discovery for 300K+ students",
    "Speech Services provide narrated content for accessibility",
    "Cosmos DB: scalable metadata store for course catalog",
]
add_multiline(sl, Inches(7.1), Inches(4.95), Inches(5.8), Inches(1.8),
              vp, sz=10, c=DARK_GRAY, bullet="→")

so_what_box(sl, Inches(0.6), Inches(5.5), Inches(6.0), Inches(1.0),
    f"At {fmtf(ai_mo)}/month, AI Course Gen represents {ai_annual/grand_total:.0%} of total plan. "
    "Pilot with 2-3 courses before full rollout to validate ROI assumptions.")

bcg_footer(sl)

# =====================================================================
# SLIDE 7 – UJIAN ONLINE RESERVED INSTANCE
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bcg_takeaway(sl,
    f"1-Year Reserved Instance commitment saves {fmtf(ri_savings)}/year across 4 regional exam clusters "
    "— each with dedicated AKS, PostgreSQL master-slave, and CDN")

# PAYG vs RI table
td = [["Service", "Qty", "PAYG/mo", "RI/mo", "Savings/mo", "Savings %"]]
for svc, desc, payg, ri, qty in RI_SERVICES:
    pt2 = payg * qty; rt = ri * qty; sv = pt2 - rt
    pct = f"{sv/pt2:.0%}" if pt2 > 0 else "–"
    td.append([svc, str(qty), fmtf(pt2), fmtf(rt), fmtf(sv), pct])
td.append(["PER REGION", "", fmtf(ri_payg_r), fmtf(ri_ri_r),
           fmtf(ri_payg_r - ri_ri_r), f"{(ri_payg_r-ri_ri_r)/ri_payg_r:.0%}"])
bcg_table(sl, td, Inches(0.6), Inches(1.2),
    [Inches(2.0), Inches(0.5), Inches(1.0), Inches(1.0), Inches(1.0), Inches(0.8)],
    hdr_bg=BCG_ORANGE)

# Region cards
add_rect(sl, Inches(7.0), Inches(1.2), Inches(6), Inches(0.35), BCG_ORANGE)
add_txt(sl, Inches(7.1), Inches(1.22), Inches(5.8), Inches(0.3),
        "4-REGION DEPLOYMENT", sz=10, bold=True, c=WHITE)

regions = [("Jakarta (JKT)", "Primary"), ("Bandung (BDG)", "Secondary"),
           ("Semarang (SMG)", "Secondary"), ("Lampung (LPG)", "Satellite")]
for i, (reg, tier) in enumerate(regions):
    x = Inches(7.1) + (i % 2) * Inches(3.0)
    y = Inches(1.7) + (i // 2) * Inches(1.4)
    add_rect(sl, x, y, Inches(2.8), Inches(1.2), WHITE, border=BCG_GREEN)
    add_rect(sl, x, y, Inches(2.8), Pt(4), BCG_GREEN)  # top accent
    add_txt(sl, x + Inches(0.1), y + Inches(0.1), Inches(2.6), Inches(0.3),
            reg, sz=12, bold=True, c=BCG_DARK_GREEN)
    add_txt(sl, x + Inches(0.1), y + Inches(0.4), Inches(2.6), Inches(0.25),
            f"RI: {fmtf(ri_ri_r)}/mo → {fmtf(ri_ri_r*12)}/yr", sz=10, c=DARK_GRAY)
    add_txt(sl, x + Inches(0.1), y + Inches(0.65), Inches(2.6), Inches(0.25),
            f"Tier: {tier}  |  Dedicated AKS", sz=8, c=LIGHT_GRAY)

# Architecture per region
add_rect(sl, Inches(0.6), Inches(4.6), Inches(12.2), Inches(0.3), BCG_DARK_GREEN)
add_txt(sl, Inches(0.7), Inches(4.62), Inches(12), Inches(0.25),
        "ARCHITECTURE PER REGION", sz=10, bold=True, c=WHITE)

arch = [
    "Vue.js Frontend (2× P1v3 App Service)", "PHP Backend – Containerized (2× D2s_v5)",
    "AKS Cluster (3× D4s_v5 nodes)", "PostgreSQL Master-Slave (D4s_v3 + D2s_v3)",
    "Azure CDN (naskah soal distribution)", "Standard Load Balancer + Log Analytics",
]
for i2, item in enumerate(arch):
    col = i2 % 3; row = i2 // 3
    add_txt(sl, Inches(0.8) + col * Inches(4.1), Inches(5.0) + row * Inches(0.28),
            Inches(3.9), Inches(0.25), f"▪  {item}", sz=9, c=DARK_GRAY)

# Summary banner
add_rect(sl, Inches(0.6), Inches(5.8), Inches(12.2), Inches(0.85), BCG_GREEN)
add_txt(sl, Inches(0.7), Inches(5.85), Inches(12), Inches(0.75),
        f"4 REGIONS TOTAL:  {fmtf(ri_ri_r*4)}/mo → {fmtf(ri_annual_4)}/year (RI)   |   "
        f"PAYG equivalent: {fmtf(ri_payg_r*4*12)}/year   |   "
        f"Annual Savings: {fmtf(ri_savings)}", sz=13, bold=True, c=WHITE, al=PP_ALIGN.CENTER)

bcg_footer(sl)

# =====================================================================
# SLIDE 8 – CONSOLIDATED BUDGET
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bcg_takeaway(sl,
    f"Total 12-month cloud investment of {fmtf(grand_total)} across three pillars, "
    f"representing a {(grand_total-fy26_annualized)/fy26_annualized:.0%} increase over FY26 annualized baseline")

# Bar chart: 3 categories + total
cd = CategoryChartData()
cd.categories = ["Existing\nSubscriptions", "AI Course\nGeneration", "Ujian Online\nRI (4 Regions)", "GRAND\nTOTAL"]
cd.add_series("Annual Cost", [plan_annual, ai_annual, ri_annual_4, grand_total])
cf = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.6), Inches(1.2), Inches(6.5), Inches(4.5), cd)
ch = cf.chart; ch.has_legend = False; ch.style = 2
pl = ch.plots[0]; pl.gap_width = 80
colors = [BCG_GREEN, BCG_TEAL, BCG_ORANGE, BCG_DARK_GREEN]
for i2, pt in enumerate(pl.series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = colors[i2]
sr = pl.series[0]; sr.has_data_labels = True
dl = sr.data_labels; dl.font.size = Pt(11); dl.font.bold = True
dl.number_format = '$#,##0'; dl.show_value = True; dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
ch.value_axis.major_gridlines.format.line.color.rgb = VERY_LIGHT_GRAY
ch.value_axis.format.line.fill.background()
ch.value_axis.tick_labels.font.size = Pt(8); ch.value_axis.tick_labels.number_format = '$#,##0'
ch.category_axis.tick_labels.font.size = Pt(9)

# Bridge / waterfall table
add_rect(sl, Inches(7.5), Inches(1.2), Inches(5.5), Inches(0.35), BCG_GREEN)
add_txt(sl, Inches(7.6), Inches(1.22), Inches(5.3), Inches(0.3),
        "INVESTMENT BUILD-UP", sz=10, bold=True, c=WHITE)

bridge = [
    ["Component", "Amount (USD)"],
    ["FY26 Annualized Baseline", fmtf(fy26_annualized)],
    ["+ Organic Growth (20%)", f"+{fmtf(plan_annual - fy26_annualized)}"],
    ["= Existing Subscriptions", fmtf(plan_annual)],
    ["+ AI Course Generation", f"+{fmtf(ai_annual)}"],
    ["+ Ujian Online RI (4 regions)", f"+{fmtf(ri_annual_4)}"],
    ["= GRAND TOTAL", fmtf(grand_total)],
    ["vs FY26 Annualized", f"+{(grand_total-fy26_annualized)/fy26_annualized:.0%}"],
]
bcg_table(sl, bridge, Inches(7.5), Inches(1.6),
    [Inches(3.0), Inches(2.0)], hdr_bg=BCG_GREEN)

# Composition bar
add_txt(sl, Inches(0.6), Inches(5.9), Inches(12), Inches(0.3),
        "BUDGET COMPOSITION", sz=10, bold=True, c=BCG_DARK_GREEN)
bar_w = Inches(12.2)
segs = [
    (plan_annual/grand_total, BCG_GREEN, f"Existing {plan_annual/grand_total:.0%}"),
    (ai_annual/grand_total, BCG_TEAL, f"AI {ai_annual/grand_total:.0%}"),
    (ri_annual_4/grand_total, BCG_ORANGE, f"RI {ri_annual_4/grand_total:.0%}"),
]
x = Inches(0.6)
for pct, clr, label in segs:
    seg_w = int(bar_w * pct)
    add_rect(sl, x, Inches(6.2), seg_w, Inches(0.45), clr)
    if pct > 0.08:
        add_txt(sl, x + Inches(0.1), Inches(6.25), seg_w - Inches(0.2), Inches(0.35),
                label, sz=9, bold=True, c=WHITE, al=PP_ALIGN.CENTER)
    x += seg_w

bcg_footer(sl)

# =====================================================================
# SLIDE 9 – KEY ASSUMPTIONS & NEXT STEPS
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bcg_takeaway(sl,
    "Execution requires phased deployment across 90 days — starting with RI procurement "
    "and pilot clusters, followed by full regional rollout")

# Left: Assumptions
add_rect(sl, Inches(0.6), Inches(1.2), Inches(5.8), Inches(0.35), BCG_BLUE)
add_txt(sl, Inches(0.7), Inches(1.22), Inches(5.6), Inches(0.3),
        "KEY ASSUMPTIONS", sz=11, bold=True, c=WHITE)

assumptions = [
    "FY26 data: Jul 2025 – Apr 2026 (10 months, April partial)",
    "20% growth applied per-month vs same FY26 calendar month",
    "Apr-26 from FY26 Apr; May/Jun use FY26 monthly average",
    "RI commitment: 1-year term, ~35% savings vs PAYG",
    "4 independent AKS clusters (JKT, BDG, SMG, Lampung)",
    "PostgreSQL master-slave per region for HA/DR",
    "All prices USD; subject to Azure pricing changes",
    "CDN shared origin Jakarta; edge POPs in other regions",
]
add_multiline(sl, Inches(0.7), Inches(1.7), Inches(5.6), Inches(3.5),
              assumptions, sz=10, c=DARK_GRAY)

# Right: Phased next steps
add_rect(sl, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.35), BCG_GREEN)
add_txt(sl, Inches(6.9), Inches(1.22), Inches(5.8), Inches(0.3),
        "IMPLEMENTATION ROADMAP", sz=11, bold=True, c=WHITE)

phases = [
    ("PHASE 1: IMMEDIATE (Week 1–2)", BCG_GREEN, [
        "Review & approve plan with CTO / CFO",
        "Procure 1-Year Reserved Instances",
        "Configure Azure Cost Management budgets & alerts",
    ]),
    ("PHASE 2: PILOT (Day 15–60)", BCG_BLUE, [
        "Deploy AKS clusters in Jakarta & Bandung",
        "Provision Azure OpenAI & AI Search",
        "Configure CDN origin + edge POPs",
    ]),
    ("PHASE 3: SCALE (Day 60–90)", BCG_TEAL, [
        "Expand AKS to Semarang & Lampung",
        "Deploy PostgreSQL master-slave per region",
        "Establish monthly consumption review cadence",
    ]),
]
y = Inches(1.7)
for phase_title, clr, items in phases:
    add_rect(sl, Inches(6.9), y, Inches(5.8), Inches(0.28), clr)
    add_txt(sl, Inches(7.0), y + Inches(0.02), Inches(5.6), Inches(0.24),
            phase_title, sz=9, bold=True, c=WHITE)
    y += Inches(0.35)
    for item in items:
        add_txt(sl, Inches(7.1), y, Inches(5.6), Inches(0.22),
                f"▪  {item}", sz=10, c=DARK_GRAY)
        y += Inches(0.25)
    y += Inches(0.1)

# Traffic light status
add_rect(sl, Inches(0.6), Inches(5.0), Inches(5.8), Inches(0.3), BCG_DARK_GREEN)
add_txt(sl, Inches(0.7), Inches(5.02), Inches(5.6), Inches(0.25),
        "READINESS ASSESSMENT", sz=10, bold=True, c=WHITE)

tl_items = [
    (BCG_GREEN, "Budget approval process", "Ready"),
    (BCG_GREEN, "Azure subscription capacity", "Ready"),
    (AMBER, "AKS cluster deployment (4 regions)", "In planning"),
    (AMBER, "AI model provisioning (GPT-4o)", "In planning"),
    (BCG_GREEN, "PostgreSQL expertise", "Available"),
    (AMBER, "CDN configuration", "In planning"),
]
for i, (clr, desc, status) in enumerate(tl_items):
    col = i % 2; row = i // 2
    x2 = Inches(0.7) + col * Inches(2.9)
    y2 = Inches(5.4) + row * Inches(0.25)
    add_rect(sl, x2, y2 + Inches(0.03), Inches(0.12), Inches(0.12), clr)
    add_txt(sl, x2 + Inches(0.18), y2, Inches(2.5), Inches(0.2),
            f"{desc} — {status}", sz=8, c=DARK_GRAY)

# Bottom banner
add_rect(sl, Inches(0.6), Inches(6.15), Inches(12.2), Inches(0.7), BCG_GREEN)
add_txt(sl, Inches(0.7), Inches(6.2), Inches(12), Inches(0.6),
        f"TOTAL INVESTMENT: {fmtf(grand_total)}  (April 2026 – March 2027)\n"
        f"Existing: {fmtf(plan_annual)}  |  AI Course Gen: {fmtf(ai_annual)}  |  "
        f"Ujian Online RI: {fmtf(ri_annual_4)}  |  RI Savings: {fmtf(ri_savings)}",
        sz=13, bold=True, c=WHITE, al=PP_ALIGN.CENTER)

bcg_footer(sl, "Source: Azure ACR Data FY26; Azure Pricing Calculator; Internal Planning Estimates")

# ── Save ──────────────────────────────────────────────────────────────
out = r"c:\Users\teddysudewo\OneDrive - Microsoft\Desktop\AccelerateDevGHCopilot\consumption-plan-ut\Azure_Consumption_Plan_UT_Apr26-Mar27_BCG.pptx"
prs.save(out)
print(f"✅ PPT saved: {out}")
print(f"   Slides: {slide_num[0]} (BCG format)")
print(f"   Grand Total: {fmtf(grand_total)}")
